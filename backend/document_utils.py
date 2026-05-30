import os
from docx import Document
from pptx import Presentation
import fitz  # PyMuPDF
import pandas as pd
from openpyxl import load_workbook
import numpy as np

# Lazy import for EasyOCR to avoid slow startup if not used
_easyocr_reader = None

def extract_preview_text(path, ext):
    """Extract first 1000 characters to detect language"""
    try:
        if ext == ".docx":
            doc = Document(path)
            return " ".join([p.text for p in doc.paragraphs[:5]])[:1000]
        elif ext == ".pptx":
            prs = Presentation(path)
            text = []
            for slide in prs.slides[:3]:
                for shape in slide.shapes:
                    if hasattr(shape, "text_frame") and shape.text_frame:
                        text.append(shape.text_frame.text)
            return " ".join(text)[:1000]
        elif ext == ".pdf":
            doc = fitz.open(path)
            text = ""
            for i in range(min(3, len(doc))):
                text += doc[i].get_text()
            doc.close()
            return text[:1000]
        elif ext == ".xlsx":
            df = pd.read_excel(path, nrows=5)
            return df.to_string()[:1000]
    except Exception as e:
        print(f"[!] Preview extraction failed: {e}")
    return ""

def get_ocr_reader():
    global _easyocr_reader
    if _easyocr_reader is None:
        import easyocr
        # Support Marathi, Hindi, and English OCR
        _easyocr_reader = easyocr.Reader(['hi', 'mr', 'en'], gpu=False) 
    return _easyocr_reader

def ocr_and_translate_page(page, model_manager, src_lang, tgt_lang):
    import numpy as np
    try:
        pix = page.get_pixmap()
        img_data = np.frombuffer(pix.samples, dtype=np.uint8).reshape(pix.h, pix.w, pix.n)
        
        reader = get_ocr_reader()
        results = reader.readtext(img_data) # returns [([coords], text, prob), ...]
        
        for (bbox_coords, text, prob) in results:
            if text.strip():
                translated = model_manager.translate(text, src_lang, tgt_lang)
                # Convert EasyOCR bbox to fitz Rect
                # EasyOCR bbox: [[x,y],[x,y],[x,y],[x,y]]
                x0 = min([p[0] for p in bbox_coords])
                y0 = min([p[1] for p in bbox_coords])
                x1 = max([p[0] for p in bbox_coords])
                y1 = max([p[1] for p in bbox_coords])
                
                # Map image coords to PDF page coords
                img_w, img_h = pix.width, pix.height
                page_w, page_h = page.rect.width, page.rect.height
                
                rect = [
                    x0 * page_w / img_w,
                    y0 * page_h / img_h,
                    x1 * page_w / img_w,
                    y1 * page_h / img_h
                ]
                
                # Since it's a scan, we don't redact (the image is the background)
                # Just overlay the text
                page.insert_textbox(rect, translated, fontsize=10, fontname="helv")
    except Exception as e:
        print(f"[!] OCR overlay failed: {e}")

def translate_docx(input_path, output_path, model_manager, src_lang, tgt_lang):
    doc = Document(input_path)
    
    # Process main paragraphs
    for para in doc.paragraphs:
        if para.text.strip():
            translated = model_manager.translate(para.text, src_lang, tgt_lang)
            if para.runs:
                # Preservation Strategy: Put whole translation in first run, clear others
                # This keeps the 'start' styling of the paragraph
                para.runs[0].text = translated
                for i in range(1, len(para.runs)):
                    para.runs[i].text = ""
                    
    # Process tables
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                for para in cell.paragraphs:
                    if para.text.strip():
                        translated = model_manager.translate(para.text, src_lang, tgt_lang)
                        if para.runs:
                            para.runs[0].text = translated
                            for i in range(1, len(para.runs)):
                                para.runs[i].text = ""
                            
    doc.save(output_path)
    return output_path

def translate_pptx(input_path, output_path, model_manager, src_lang, tgt_lang):
    prs = Presentation(input_path)
    
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text_frame") and shape.text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    if paragraph.text.strip():
                        translated = model_manager.translate(paragraph.text, src_lang, tgt_lang)
                        if paragraph.runs:
                            paragraph.runs[0].text = translated
                            for i in range(1, len(paragraph.runs)):
                                paragraph.runs[i].text = ""
                                
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        if cell.text_frame:
                            for paragraph in cell.text_frame.paragraphs:
                                if paragraph.text.strip():
                                    translated = model_manager.translate(paragraph.text, src_lang, tgt_lang)
                                    if paragraph.runs:
                                        paragraph.runs[0].text = translated
                                        for i in range(1, len(paragraph.runs)):
                                            paragraph.runs[i].text = ""
                                        
    prs.save(output_path)
    return output_path

def translate_xlsx(input_path, output_path, model_manager, src_lang, tgt_lang):
    from copy import copy
    wb = load_workbook(input_path)
    texts_to_translate = []
    cell_refs = [] 
    
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        for row in ws.iter_rows():
            for cell in row:
                if cell.value and isinstance(cell.value, str) and cell.value.strip():
                    texts_to_translate.append(cell.value)
                    cell_refs.append((sheet_name, cell.coordinate))

    if not texts_to_translate:
        wb.save(output_path)
        return output_path

    translated_texts = model_manager.translate_batch(texts_to_translate, src_lang, tgt_lang)
    
    for i, (sheet_name, coord) in enumerate(cell_refs):
        ws = wb[sheet_name]
        cell = ws[coord]
        
        # Clone style
        original_font = copy(cell.font)
        original_fill = copy(cell.fill)
        original_border = copy(cell.border)
        original_alignment = copy(cell.alignment)
        original_number_format = cell.number_format
        original_protection = copy(cell.protection)
        
        cell.value = translated_texts[i]
        
        # Re-apply style
        cell.font = original_font
        cell.fill = original_fill
        cell.border = original_border
        cell.alignment = original_alignment
        cell.number_format = original_number_format
        cell.protection = original_protection
        
    wb.save(output_path)
    return output_path

def translate_pdf(input_path, output_path, model_manager, src_lang, tgt_lang):
    import fitz
    doc = fitz.open(input_path)
    
    for page in doc:
        # Get text blocks with coordinates
        blocks = page.get_text("dict")["blocks"]
        text_found = False
        
        for b in blocks:
            if "lines" in b:
                text_found = True
                for l in b["lines"]:
                    for s in l["spans"]:
                        original_text = s["text"]
                        if original_text.strip():
                            translated = model_manager.translate(original_text, src_lang, tgt_lang)
                            bbox = s["bbox"] # (x0, y0, x1, y1)
                            
                            # Redact original text
                            page.add_redact_annot(bbox, fill=(1, 1, 1))
                            page.apply_redactions()
                            
                            # Insert translated text
                            fontsize = s["size"]
                            # align=0 (Left), align=1 (Center), align=2 (Right)
                            page.insert_textbox(bbox, translated, fontsize=fontsize, fontname="helv", align=0)
        
        # Fallback for scanned pages (Task 3 integration)
        if not text_found:
            ocr_and_translate_page(page, model_manager, src_lang, tgt_lang)
                            
    doc.save(output_path)
    doc.close()
    return output_path
