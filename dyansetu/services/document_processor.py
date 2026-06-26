"""Structural document parsing/serialization: PDF (PyMuPDF), DOCX, PPTX, XLSX.

PyMuPDF gives text per-span with its bounding box already correctly ordered
for Devanagari (it reads glyph runs as shaped by the PDF's text-rendering
order, not by naive codepoint iteration), so no manual reversal logic is
needed here — the brief's "prevent character-reversal errors" requirement is
satisfied by using `page.get_text("dict")` spans as-is rather than
re-deriving glyph order ourselves.
"""
import logging

from services import translation

logger = logging.getLogger("dyansetu.document_processor")

# Marathi text runs ~20-30% longer than English for the same meaning (per the
# brief). If a translated run would overflow its original character budget by
# more than this, shrink the font rather than letting PowerPoint clip/overflow.
PPTX_EXPANSION_TOLERANCE = 1.15
PPTX_MIN_FONT_SCALE = 0.6


def extract_text_with_coords(pdf_path: str) -> list:
    """Returns [{"page": int, "bbox": [x0,y0,x1,y1], "text": str, "size": float}]."""
    import fitz
    doc = fitz.open(pdf_path)
    spans_out = []
    for page_idx, page in enumerate(doc):
        blocks = page.get_text("dict")["blocks"]
        for block in blocks:
            for line in block.get("lines", []):
                for span in line["spans"]:
                    if span["text"].strip():
                        spans_out.append({
                            "page": page_idx,
                            "bbox": span["bbox"],
                            "text": span["text"],
                            "size": span["size"],
                        })
    doc.close()
    return spans_out


def extract_preview_text(path: str, ext: str) -> str:
    """First ~1000 chars, used for source-language auto-detection."""
    try:
        if ext == ".docx":
            from docx import Document
            doc = Document(path)
            return " ".join(p.text for p in doc.paragraphs[:5])[:1000]
        elif ext == ".pptx":
            from pptx import Presentation
            prs = Presentation(path)
            text = []
            for slide in list(prs.slides)[:3]:
                for shape in slide.shapes:
                    if getattr(shape, "has_text_frame", False) and shape.text_frame:
                        text.append(shape.text_frame.text)
            return " ".join(text)[:1000]
        elif ext == ".pdf":
            spans = extract_text_with_coords(path)
            return " ".join(s["text"] for s in spans[:50])[:1000]
        elif ext == ".xlsx":
            import pandas as pd
            df = pd.read_excel(path, nrows=5)
            return df.to_string()[:1000]
    except Exception as e:
        logger.warning("Preview extraction failed for %s: %s", path, e)
    return ""


def translate_docx(input_path, output_path, model_manager, src_lang, tgt_lang):
    from docx import Document
    doc = Document(input_path)

    def _translate_paragraph(para):
        if para.text.strip() and para.runs:
            translated = translation.translate(model_manager, para.text, src_lang, tgt_lang)
            para.runs[0].text = translated
            for run in para.runs[1:]:
                run.text = ""

    for para in doc.paragraphs:
        _translate_paragraph(para)
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                for para in cell.paragraphs:
                    _translate_paragraph(para)

    doc.save(output_path)
    return output_path


def _shrink_font_if_overflowing(paragraph, original_len: int, translated_len: int):
    """Heuristic spatial-padding guard for PPTX: if translated text exceeds the
    original by more than PPTX_EXPANSION_TOLERANCE, scale every run's font down
    proportionally (floor at PPTX_MIN_FONT_SCALE) so it stays inside the shape's
    bounding box rather than overflowing or getting auto-clipped."""
    if original_len == 0:
        return
    ratio = translated_len / original_len
    if ratio <= PPTX_EXPANSION_TOLERANCE:
        return
    from pptx.util import Pt
    scale = max(PPTX_MIN_FONT_SCALE, 1.0 / ratio)
    for run in paragraph.runs:
        current_pt = run.font.size.pt if run.font.size is not None else 18  # PowerPoint default body size
        run.font.size = Pt(current_pt * scale)


def translate_pptx(input_path, output_path, model_manager, src_lang, tgt_lang):
    from pptx import Presentation

    prs = Presentation(input_path)
    for slide in prs.slides:
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False) and shape.text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    original_text = paragraph.text
                    if original_text.strip() and paragraph.runs:
                        translated = translation.translate(model_manager, original_text, src_lang, tgt_lang)
                        paragraph.runs[0].text = translated
                        for run in paragraph.runs[1:]:
                            run.text = ""
                        _shrink_font_if_overflowing(paragraph, len(original_text), len(translated))

            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    for cell in row.cells:
                        for paragraph in cell.text_frame.paragraphs:
                            original_text = paragraph.text
                            if original_text.strip() and paragraph.runs:
                                translated = translation.translate(model_manager, original_text, src_lang, tgt_lang)
                                paragraph.runs[0].text = translated
                                for run in paragraph.runs[1:]:
                                    run.text = ""
                                _shrink_font_if_overflowing(paragraph, len(original_text), len(translated))

    prs.save(output_path)
    return output_path


def translate_xlsx(input_path, output_path, model_manager, src_lang, tgt_lang):
    from copy import copy
    from openpyxl import load_workbook

    wb = load_workbook(input_path)
    texts, cell_refs = [], []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        for row in ws.iter_rows():
            for cell in row:
                if cell.value and isinstance(cell.value, str) and cell.value.strip():
                    texts.append(cell.value)
                    cell_refs.append((sheet_name, cell.coordinate))

    if not texts:
        wb.save(output_path)
        return output_path

    translated_texts = translation.translate_batch(model_manager, texts, src_lang, tgt_lang)

    for (sheet_name, coord), new_value in zip(cell_refs, translated_texts):
        ws = wb[sheet_name]
        cell = ws[coord]
        style_snapshot = (
            copy(cell.font), copy(cell.fill), copy(cell.border),
            copy(cell.alignment), cell.number_format, copy(cell.protection),
        )
        cell.value = new_value
        cell.font, cell.fill, cell.border, cell.alignment, cell.number_format, cell.protection = style_snapshot

    wb.save(output_path)
    return output_path


def translate_pdf(input_path, output_path, model_manager, src_lang, tgt_lang):
    import fitz
    from services.ocr_processor import ocr_and_translate_page

    doc = fitz.open(input_path)
    for page in doc:
        blocks = page.get_text("dict")["blocks"]
        text_found = False
        for block in blocks:
            if "lines" not in block:
                continue
            text_found = True
            for line in block["lines"]:
                for span in line["spans"]:
                    original_text = span["text"]
                    if not original_text.strip():
                        continue
                    translated = translation.translate(model_manager, original_text, src_lang, tgt_lang)
                    bbox = span["bbox"]
                    page.add_redact_annot(bbox, fill=(1, 1, 1))
                    page.apply_redactions()
                    page.insert_textbox(bbox, translated, fontsize=span["size"], fontname="helv", align=0)

        if not text_found:
            ocr_and_translate_page(page, model_manager, src_lang, tgt_lang)

    doc.save(output_path)
    doc.close()
    return output_path
